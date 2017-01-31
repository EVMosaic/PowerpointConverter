import os
import tkinter
from tkinter import filedialog
from tkinter import constants
from pptx import Presentation
from mappings import templates, mappings, SlideType

class Application:
    save_path = ''
    ppt_path = ''
    root = tkinter.Tk()
    powerpoint = None
    status = tkinter.StringVar()
    js_update = None

    def __init__(self):
        self.build_gui()
        Application.root.mainloop()

    def build_gui(self):
        self.root.minsize(300,100)
        self.root.lift()
        button_opts = {'fill': constants.BOTH, 'padx': 5, 'pady': 5}
        ppt_button = tkinter.Button(Application.root, text='Select Powerpoint', command=self.set_ppt_path).pack(**button_opts)
        save_button = tkinter.Button(Application.root, text='Set Save Location', command=self.set_save_path).pack(**button_opts)
        start_button = tkinter.Button(Application.root, text='Convert', command=self.run).pack(**button_opts)
        label = tkinter.Label(Application.root, textvariable=Application.status).pack()
        self.update_status('Please Select Powerpoint')

    def set_ppt_path(self):
        opts = {'parent': Application.root,
                'title': 'Select Powerpoint to Convert',
                'filetypes':  [('Powerpoint Files', '.pptx')],
                'initialdir': 'C:/Users/eric_/Desktop/GreenMockups/Production Tests'}  # change to C:/

        Application.ppt_path = filedialog.askopenfilename(**opts)
        Application.update_status('Powerpoint Loaded')

    def set_save_path(self):
        opts = {'parent':  Application.root,
                'title': 'Select Save Location',
                'initialdir': 'C:/Users/eric_/Desktop/GreenMockups/eta-sample/eta-sample/pages'}  # change to C:/


        Application.save_path = filedialog.askdirectory(**opts)
        Application.update_status('Output Set')

        Application.js_update = JsUpdater(Application.save_path)


    def run(self):
        print('Starting application')
        Application.status.set('')
        Application.powerpoint = Powerpoint(Application.ppt_path)
        Converter.convert_presentation(Application.powerpoint, templates, mappings)
        Application.js_update.save_to_disk()

    @staticmethod
    def update_status(new_status):
        Application.status.set(Application.status.get() + '\n' + new_status)

class Powerpoint:
    def __init__(self, ppt_path):
        self.powerpoint = Presentation(ppt_path)
        self.layouts = []
        self.names = []
        self.make_layout_list()
        self.name_slides()

    def make_layout_list(self):
        for i in range(len(self.powerpoint.slide_layouts)):
            self.layouts.append(self.powerpoint.slide_layouts[i])

    def name_slides(self):
        for slide in self.powerpoint.slides:
            slide_number = self.powerpoint.slides.index(slide) + 1
            slide.name = "Slide%02d" % slide_number
            self.names.append(slide.name)

    def get_slide_type(self, slide):
        return SlideType(self.layouts.index(slide.slide_layout))


class Page:
    template_head = '''<!DOCTYPE html>
    <html lang="en">
    <head>
        <meta charset="utf-8">
        <meta http-equiv="X-UA-Compatible" content="IE=edge">
        <meta http-equiv="Cache-Control" Content="no-cache" />
        <meta http-equiv="Pragma" Content="no-cache" />
        <meta http-equiv="Expires" Content="0" />
        <meta name="viewport" content="width=device-width, initial-scale=1">
        <meta name="apple-mobile-web-app-capable" content="yes">
        <meta name="apple-mobile-web-app-status-bar-style" content="black">
        <meta name="format-detection" content="telephone=no">

        <title></title>

        <link href="https://fonts.googleapis.com/css?family=Titillium+Web:300,400,400i,700" rel="stylesheet">
        <link href="https://fonts.googleapis.com/css?family=Permanent+Marker" rel="stylesheet">
        <link href="https://fonts.googleapis.com/css?family=Caveat:400,700" rel="stylesheet">
        <link href="https://fonts.googleapis.com/css?family=Open+Sans:400,400i,700,700i" rel="stylesheet">
        <link href="https://maxcdn.bootstrapcdn.com/font-awesome/4.7.0/css/font-awesome.min.css" rel="stylesheet">
        <link href="../../css/animate.min.css" rel="stylesheet" >
        <link href="../../css/sandstone.css" rel="stylesheet">
        <link href="../../css/main.css" rel="stylesheet">
        <link rel="apple-touch-icon-precomposed" href="../../images/_icons/touch-icon-iphone-60.png">
        <link rel="apple-touch-icon-precomposed" sizes="76x76" href="../../images/_icons/touch-icon-ipad-76.png">
        <link rel="apple-touch-icon-precomposed" sizes="120x120" href="../../images/_icons/touch-icon-iphone-retina-120.png">
        <link rel="apple-touch-icon-precomposed" sizes="152x152" href="../../images/_icons/touch-icon-ipad-retina-152.png">
        <!-- HTML5 shim and Respond.js for IE8 support of HTML5 elements and media queries -->
        <!-- WARNING: Respond.js doesn't work if you view the page via file:// -->
        <!--[if lt IE 9]>
          <script src="https://oss.maxcdn.com/html5shiv/3.7.3/html5shiv.min.js"></script>
          <script src="https://oss.maxcdn.com/respond/1.4.2/respond.min.js"></script>
        <![endif]-->
    </head>
    <body id="main">

        <!--==================================================
        ======================================================
        START CONTENT - EDIT BELOW ONLY
        ======================================================
        ===================================================-->


        <!--===============================================================
        Section
        ================================================================-->
        '''
    template_tail = '''
    <!--==================================================
        ======================================================
        END CONTENT - EDIT ABOVE ONLY
        ======================================================
        ===================================================-->



        <script src="../../js/jquery-3.1.1.min.js"></script>
        <script src="../../js/bootstrap.min.js"></script>
        <script src="../../js/iframeResizer.contentWindow.min.js" defer></script>
        <script src="../../js/page.js"></script>


    </body>
    </html>'''

    def __init__(self, name):
        self.html = ''
        self.images = {}
        self.slide_name = name

    def build_html(self, html):
        self.html = Page.template_head + html + Page.template_tail

    def add_image(self, name, image):
        self.images[name] = image

    def save_to_disk(self):
        base_path = os.path.join(Application.save_path,  self.slide_name)
        media_path = os.path.join(base_path, 'media')
        page_name = os.path.join(base_path,  'index.html')

        os.makedirs(base_path, exist_ok=True)
        os.makedirs(media_path, exist_ok=True)

        with open(page_name, 'w', encoding='utf8') as new_file:
            new_file.write(self.html)

        for image_name in self.images:
            image_path = os.path.join(media_path, image_name)
            with open(image_path, 'wb') as image:
                image.write(self.images[image_name])


class Converter:
    @staticmethod
    def convert_slide(slide, template, mapping):
        current_task = 'Starting conversion of ' + slide.name;
        Application.update_status(current_task)
        print('Starting conversion of ', slide.name)
        html = template
        mapping = mapping
        page = Page(slide.name)
        image_count = 0

        for item in mapping:
            #print('Looking up item.idx: %s item.template_element: %s, item.slide_element: %s' % (item.idx, item.template_element, item.slide_element))
            try:
                element = slide.placeholders[item.idx]
                if element.has_text_frame:
                    formatted_text = Converter.make_tags(Converter.sanitize(element.text), 'p') + '\n'
                    print(formatted_text)
                    html = html.replace(item.template_element, formatted_text)
                else:  # for the time being can reasonably assume this means an image
                    # but should probably figure out a better way to handle this
                    # this is already broken since tables fail this check
                    img = element.image
                    filename = 'image%d.%s' % (image_count, img.ext)
                    page.add_image(filename, img.blob)
                    img_src = '"media/%s"' % filename
                    print('adding image at ', img_src, ' replacing element ', item.template_element)
                    html = html.replace(item.template_element, img_src)
                    image_count += 1
            except:
                print('encountered error on ', slide.name)
                Application.update_status('[ ENCOUNTERED ERROR ON ' + slide.name.upper() + ' ]')
                error_msg = '</IMG> <h1 style="color:white; background-color:red; padding:20px; text-align:center; font-size:40px; margin:0">' + 'IMAGE OR TEXT NEEDS MANUAL REPLACEMENT' + '</h1>' +  '<h2 style="color:yellow; background-color:red; padding-bottom:20px; text-align:center; font-size:25px; margin:0">' + 'Encountered error on - Slide Element:' + str(item.slide_element) + ' with IDX:' + str(item.idx) + ' for Template Element ' + str(item.template_element) + '</h2>'

                html = html.replace(item.template_element, error_msg)

        page.build_html(html)
        return page

    @staticmethod
    def convert_presentation(powerpoint, templates, mappings):
        for slide in powerpoint.powerpoint.slides:
            slide_type = powerpoint.get_slide_type(slide)
            template = templates[slide_type]
            mapping = mappings[slide_type]
            page = Converter.convert_slide(slide, template, mapping)
            page.save_to_disk()
        Application.update_status('Finished Conversion')

    @staticmethod
    def make_tags(text, tag):
        return '<%s>' % tag + text + '</%s>' % tag

    @staticmethod
    def sanitize(text):
        sani_text = text
        replacements = {
                        }
        for dirty, clean in replacements.items():
            sani_text = sani_text.replace(dirty, clean)

        return sani_text

class JsUpdater:

    def __init__(self, save_path):
        self.path = os.path.normpath(os.path.join(save_path, os.path.pardir, 'js', 'main.js' ))

    def make_comment(self, comment):
        return "/*==================================================\n" + comment + "\n==================================================*/\n"

    def make_pages_list(self, pages):
        pages_list = "const pages = ["
        for page in pages:
            pages_list += '\n    "pages/{}/index.html,"'.format(page)
        pages_list += "\n];\n"
        return pages_list

    def save_to_disk(self):
        main_js = open(self.path)
        template = main_js.read()
        main_js.close()
        start = template.find('const pages = [')
        end = template.find(']', start) + 3
        updated = template.replace(template[start:end], self.make_pages_list(Application.powerpoint.names))

        with open(self.path, 'w', encoding='utf8') as main_js:
            main_js.write(updated)


if __name__ == '__main__':
    app = Application()
